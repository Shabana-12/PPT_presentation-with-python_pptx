# Modules to be imported
from pptx import Presentation	
import os
from datetime import date
from pptx.util import Inches,Pt,Cm
from wand.image import Image
from wand.compat import nested
from wand.color import Color
files = os.listdir("D:\\Home Python Assignment")

# ny = Image(filename ='nike_black.png')
# print(ny.width, ny.height)
# ny.resize(1500,500)
# print(ny.width, ny.height)
# ny.save(filename ='logo.png')
root = Presentation()
slide_layout = root.slide_layouts[6]

def add_slide(root, layout,image_url):
    slide = root.slides.add_slide(layout)
    
    txBox = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(20), Cm(2))
    title= txBox.text_frame
    title.text = "My presentation with python-pptx"
    subtitle= slide.shapes.add_textbox(Cm(2), Cm(2), Cm(15), Cm(1))
    # set font 
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())
    left = top = Inches(1.5) 
   
    slide.shapes.add_picture(image_url, left, top, Inches(3.5))
    
    return slide

img_files = list(filter(lambda i: '.jpg' in i, files))
for i in img_files:
    with Image(filename=i) as background:
        with Image(filename='logo.png') as watermark:
            background.watermark(image=watermark, transparency=0.33)
            background.save(filename=i)
        add_slide(root,slide_layout,i)

root.save("PPT.pptx")
os.startfile("PPT.pptx")
print("Presentation Done Successfully")